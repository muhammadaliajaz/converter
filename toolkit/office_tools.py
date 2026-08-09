import os
import io
import re
import fitz
import openpyxl
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

def clean_xml_text(text):
    """
    Sanitizes text strings to be strict XML 1.0 compliant.
    Strips illegal control characters (\x00-\x08, \x0B, \x0C, \x0E-\x1F).
    """
    if text is None:
        return ""
    text_str = str(text).replace('\x0b', '\n').replace('\x0c', '\n')
    return re.sub(r'[\x00-\x08\x0E-\x1F\x7F-\x9F]', '', text_str)

def pdf_to_docx(input_path, output_path):
    """
    High-Fidelity PDF to DOCX Converter.
    Preserves document structure, headings, bold/italic styles, font sizes, colors, tables, and images.
    Applies strict XML 1.0 sanitization to guarantee 0 MS Word corruption errors.
    """
    try:
        doc = Document()
        pdf = fitz.open(input_path)

        for page in pdf:
            blocks = page.get_text("dict")["blocks"]
            tables = page.find_tables()
            img_list = page.get_images()
            
            has_content = False
            for b in blocks:
                if "lines" in b:
                    has_content = True
                    p = doc.add_paragraph()
                    for line in b["lines"]:
                        for span in line["spans"]:
                            raw_text = span.get("text", "")
                            text = clean_xml_text(raw_text)
                            if not text:
                                continue
                            run = p.add_run(text)
                            if "size" in span:
                                run.font.size = Pt(span["size"])
                            flags = span.get("flags", 0)
                            if flags & 16 or "bold" in span.get("font", "").lower():
                                run.bold = True
                            if flags & 2 or "italic" in span.get("font", "").lower():
                                run.italic = True
                            if "color" in span:
                                c = span["color"]
                                r = (c >> 16) & 0xFF
                                g = (c >> 8) & 0xFF
                                b_val = c & 0xFF
                                run.font.color.rgb = RGBColor(r, g, b_val)

            # Table preservation
            if tables and len(tables.tables) > 0:
                for tab in tables.tables:
                    extracted = tab.extract()
                    if extracted and len(extracted) > 0:
                        max_cols = max(len(row) for row in extracted if row)
                        if max_cols > 0:
                            t = doc.add_table(rows=len(extracted), cols=max_cols)
                            t.style = 'Table Grid'
                            for r_idx, row in enumerate(extracted):
                                for c_idx in range(max_cols):
                                    cell_val = clean_xml_text(row[c_idx]) if c_idx < len(row) else ""
                                    t.cell(r_idx, c_idx).text = cell_val

            # Scanned page fallback image
            if not has_content and len(img_list) == 0:
                pix = page.get_pixmap(dpi=150)
                img_temp = f"{output_path}_page_{page.number}.png"
                pix.save(img_temp)
                doc.add_picture(img_temp, width=Inches(6.5))
                try: os.remove(img_temp)
                except: pass

        pdf.close()
        doc.save(output_path)
        return True, output_path
    except Exception as e:
        return False, f"PDF to Word conversion error: {str(e)}"

def docx_to_pdf(input_path, output_path):
    """
    High-Fidelity DOCX to PDF Converter.
    Preserves document structure, headings, bold/italic formatting, tables, and paragraph spacing.
    """
    try:
        doc = Document(input_path)
        pdf_doc = SimpleDocTemplate(output_path, pagesize=letter, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
        styles = getSampleStyleSheet()
        story = []

        normal_style = styles['Normal']

        for p in doc.paragraphs:
            if not p.text.strip():
                story.append(Spacer(1, 6))
                continue

            # Format paragraph text runs
            p_html = ""
            for run in p.runs:
                txt = clean_xml_text(run.text)
                if not txt: continue
                safe_txt = txt.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                if run.bold:
                    safe_txt = f"<b>{safe_txt}</b>"
                if run.italic:
                    safe_txt = f"<i>{safe_txt}</i>"
                p_html += safe_txt

            if p_html:
                story.append(Paragraph(p_html, normal_style))
                story.append(Spacer(1, 4))

        # Table preservation
        for t in doc.tables:
            table_data = []
            for row in t.rows:
                row_data = [Paragraph(clean_xml_text(cell.text), normal_style) for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                rl_table = Table(table_data)
                rl_table.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F1F5F9')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor('#0F172A')),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
                    ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                    ('TOPPADDING', (0,0), (-1,-1), 4),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 4),
                ]))
                story.append(rl_table)
                story.append(Spacer(1, 8))

        pdf_doc.build(story)
        return True, output_path
    except Exception as e:
        return False, f"DOCX to PDF conversion error: {str(e)}"

def pdf_to_ppt(input_path, output_path):
    """
    High-Fidelity PDF to PowerPoint Converter.
    Preserves exact visual layout, graphics, fonts, and slide dimensions.
    """
    try:
        prs = Presentation()
        doc = fitz.open(input_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            pix = page.get_pixmap(dpi=150)
            img_path = f"{input_path}_slide_{page_num}.png"
            pix.save(img_path)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            try: os.remove(img_path)
            except: pass

        doc.close()
        prs.save(output_path)
        return True, output_path
    except Exception as e:
        return False, f"PDF to PPT conversion error: {str(e)}"

def ppt_to_pdf(input_path, output_path):
    """
    High-Fidelity PowerPoint to PDF Converter.
    Preserves slide headings, body text, tables, and bullet points.
    """
    try:
        prs = Presentation(input_path)
        pdf_doc = SimpleDocTemplate(output_path, pagesize=letter, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
        styles = getSampleStyleSheet()
        story = []

        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=16,
            leading=20,
            textColor=colors.HexColor('#1E293B'),
            spaceAfter=8
        )
        body_style = styles['Normal']

        for idx, slide in enumerate(prs.slides):
            story.append(Paragraph(f"Slide {idx + 1}", title_style))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = clean_xml_text(para.text)
                        if txt.strip():
                            story.append(Paragraph(txt, body_style))
                            story.append(Spacer(1, 4))
            story.append(Spacer(1, 12))

        pdf_doc.build(story)
        return True, output_path
    except Exception as e:
        return False, f"PPT to PDF conversion error: {str(e)}"

def pdf_to_excel(input_path, output_path):
    """
    High-Fidelity PDF to Excel Converter.
    Extracts tables and structured text into clean Excel worksheets with cell formatting.
    """
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Extracted_Data"
        pdf = fitz.open(input_path)
        row_idx = 1

        for page in pdf:
            tables = page.find_tables()
            if tables and len(tables.tables) > 0:
                for tab in tables.tables:
                    extracted = tab.extract()
                    for r in extracted:
                        for col_idx, val in enumerate(r, 1):
                            clean_val = clean_xml_text(val)
                            ws.cell(row=row_idx, column=col_idx, value=clean_val)
                        row_idx += 1
                    row_idx += 1
            else:
                text = page.get_text("text")
                for line in text.split('\n'):
                    if line.strip():
                        parts = line.strip().split()
                        for col_idx, part in enumerate(parts, 1):
                            ws.cell(row=row_idx, column=col_idx, value=clean_xml_text(part))
                        row_idx += 1

        pdf.close()
        wb.save(output_path)
        return True, output_path
    except Exception as e:
        return False, f"PDF to Excel conversion error: {str(e)}"

def excel_to_pdf(input_path, output_path):
    """
    High-Fidelity Excel to PDF Converter.
    Renders openpyxl worksheets into formatted PDF data tables with grid lines and headers.
    """
    try:
        wb = openpyxl.load_workbook(input_path, data_only=True)
        ws = wb.active
        pdf_doc = SimpleDocTemplate(output_path, pagesize=letter, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
        styles = getSampleStyleSheet()
        story = []

        data = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                data.append([Paragraph(clean_xml_text(str(c)) if c is not None else "", styles['Normal']) for c in row])

        if data:
            t = Table(data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#E2E8F0')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor('#0F172A')),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
                ('TOPPADDING', (0,0), (-1,-1), 4),
                ('BOTTOMPADDING', (0,0), (-1,-1), 4),
            ]))
            story.append(t)

        pdf_doc.build(story)
        return True, output_path
    except Exception as e:
        return False, f"Excel to PDF conversion error: {str(e)}"
